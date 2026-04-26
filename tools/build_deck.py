"""
tools/build_deck.py — assemble the 3-minute hackathon deck as horizon_pitch.pptx.

Slide plan (16:9, ~3 min total):
    1. Title / hook
    2. Problem & motivation
    3. Easy dataset — bar chart + overlay video
    4. Hard dataset — bar chart + overlay video
    5. Algorithm walkthrough A — input -> Lab -> Otsu masks
    6. Algorithm walkthrough B — boundary -> RANSAC pool -> coherence rerank
    7. Performance testing — Pi 5 docker model + latency results
    8. Final result + budget recap
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

REPO_ROOT = Path(__file__).resolve().parent.parent
PRES_DIR = REPO_ROOT / "presentation"
HERO_DIR = PRES_DIR / "hero_pipeline"
CHART_DIR = PRES_DIR / "charts"
VIDEO_DIR = PRES_DIR / "videos"

OUT_PATH = PRES_DIR / "horizon_pitch.pptx"

# Colour palette
ACCENT_BLUE = RGBColor(0x1F, 0x4E, 0x8C)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x2C, 0xA0, 0x2C)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)


def set_slide_size_16x9(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_blank_slide(prs: Presentation):
    """Use the layout with the fewest placeholders so we control everything."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(34)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT_BLUE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.runs[0].font.size = Pt(16)
        sp.runs[0].font.color.rgb = TEXT_MUTED


def add_text_block(slide, left, top, width, height, lines, *, size=18, bold_first=False, color=TEXT_DARK) -> None:
    """`lines` is a list of (text, is_bold, font_size_override_or_none) or plain str."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, is_bold, override = line, (bold_first and i == 0), None
        else:
            text, is_bold, override = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Empty paragraphs have no runs; insert a space so styling is applied
        # and the line still takes vertical space.
        p.text = text if text else " "
        run = p.runs[0]
        run.font.size = Pt(override if override else size)
        run.font.bold = is_bold
        run.font.color.rgb = color
        p.space_after = Pt(6)


def add_image(slide, path: Path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_video_with_autoplay_loop(slide, video_path: Path, poster_path: Path, left, top, width, height):
    """Embed a video and force autoplay + loop via raw XML.

    python-pptx exposes add_movie() but doesn't expose the autoplay/loop knobs,
    so we patch the timing nodes after insertion. Without this, PowerPoint
    waits for a click on every slide transition — which kills the flow of a
    3-minute talk.
    """
    movie = slide.shapes.add_movie(
        str(video_path),
        left, top, width, height,
        poster_frame_image=str(poster_path),
    )

    # Step 1 — set the picture's nvPicPr/nvPr/extLst child <a:videoFile> media
    # reference to "loop". This is what makes the video keep playing after it
    # finishes; PowerPoint defaults to "play once".
    pic = movie._element  # <p:pic>
    nvPr = pic.find(qn("p:nvPicPr")).find(qn("p:nvPr"))
    # Add an <a:videoFile loop="true"> hint via the existing extLst/cMediaNode.
    for child in nvPr.iter(qn("p:videoFile")):
        # No-op: videoFile already exists; we'll patch the timeline below.
        pass

    # Step 2 — find the <p:timing> tree on the slide, locate the media node
    # that points back at this picture's id, and set @loop="true" plus the
    # play behaviour to "fromBeginning" autoplay.
    sp_id = pic.find(qn("p:nvPicPr")).find(qn("p:nvPr")).find(qn("p:nvPr"))  # placeholder
    # The shape id is on <p:cNvPr id="...">.
    cNvPr = pic.find(qn("p:nvPicPr")).find(qn("p:cNvPr"))
    shape_id = cNvPr.get("id")

    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        return movie

    # Walk timing tree, find <p:video><p:cMediaNode spid="<shape_id>">.
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    for video_node in timing.iterfind(".//p:video", nsmap):
        cMediaNode = video_node.find(qn("p:cMediaNode"))
        if cMediaNode is None:
            continue
        # Looping: cMediaNode @loop and cTn @repeatCount.
        cMediaNode.set("loop", "1")
        cTn = cMediaNode.find(qn("p:cTn"))
        if cTn is not None:
            cTn.set("repeatCount", "indefinite")
            cTn.set("fill", "hold")
            # Autoplay on slide load: the default node type is "afterEffect"
            # with delay="indefinite", which makes PowerPoint wait for a
            # click. Switch to "withEffect" + delay=0 so the video plays
            # automatically as soon as the slide is shown.
            cTn.set("nodeType", "withEffect")
            stCondLst = cTn.find(qn("p:stCondLst"))
            if stCondLst is not None:
                for cond in stCondLst.findall(qn("p:cond")):
                    cond.set("delay", "0")

    # Step 3 — autoplay flag on the videoFile element itself.
    for vfile in pic.iter(qn("p:videoFile")):
        # No standard attribute; some readers honour 'r:link' via extLst. Skip.
        pass

    return movie


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.color.rgb = TEXT_MUTED


def add_colored_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = add_blank_slide(prs)
    # Background band
    add_colored_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BG_LIGHT)
    add_colored_rect(slide, Inches(0), Inches(2.6), prs.slide_width, Inches(0.05), ACCENT_BLUE)

    # Hero title
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.4))
    p = box.text_frame.paragraphs[0]
    p.text = "Real-time horizon detection on the Pi 5"
    p.runs[0].font.size = Pt(54)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT_BLUE

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.0))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Crop the sky out of the frame before the detector sees it - within 67 ms / frame budget."
    sp.runs[0].font.size = Pt(22)
    sp.runs[0].font.color.rgb = TEXT_DARK

    # Tag line
    tag = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "10 attempts. ~30 ms on the Pi 5 model. 33 FPS while detector runs on Hailo."
    tp.runs[0].font.size = Pt(15)
    tp.runs[0].font.italic = True
    tp.runs[0].font.color.rgb = TEXT_MUTED


def slide_problem(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "The problem", "Half the frame is sky - and the sky lies to you")

    # Body bullets — left column
    bullets = [
        ("UAV at 50-60 m: half the frame or more is sky.", True, 22),
        ("The ground-target detector spends Hailo cycles on pixels that can't contain a target.", False, 18),
        ("Worse: clouds, sun glare and birds masquerade as ground targets - false positives.", False, 18),
        ("", False, 8),
        ("Fix: crop above the horizon before the detector runs.", True, 22),
        ("Catch: the horizon detector itself must be fast enough to be worth running.", False, 18),
    ]
    add_text_block(slide, Inches(0.5), Inches(1.7), Inches(7.6), Inches(4.5), bullets, size=18)

    # Hard requirements box on the right
    add_colored_rect(slide, Inches(8.4), Inches(1.7), Inches(4.5), Inches(4.7), BG_LIGHT, ACCENT_BLUE)
    req = slide.shapes.add_textbox(Inches(8.6), Inches(1.85), Inches(4.2), Inches(4.5))
    tf = req.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Hard requirements"
    p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT_BLUE
    for line in [
        "≥ 15 FPS on Pi 5 CPU",
        "  while a detector runs on Hailo",
        "",
        "Output: angle + offset",
        "  (or sky mask)",
        "",
        "Onboard - no cloud calls",
        "",
        "Internal pass gate:",
        "  Δθ < 5°  AND  Δρ/H < 5%",
    ]:
        p = tf.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = TEXT_DARK

    add_footer(slide, "Source: hackathon brief / project CLAUDE.md")


def slide_easy_dataset(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Easy dataset - Horizon-UAV (490 frames)",
              "Aerial UAV footage. Sky/ground split, sometimes hazy, mild-to-moderate roll.")

    # Bar chart (left half)
    add_image(slide, CHART_DIR / "uav_progression.png",
              Inches(0.4), Inches(1.6), width=Inches(7.4))

    # Video (right half)
    add_video_with_autoplay_loop(
        slide,
        VIDEO_DIR / "uav_overlay.mp4",
        HERO_DIR / "07_final.png",  # poster frame: the final hero output
        Inches(8.1), Inches(1.7), Inches(4.9), Inches(4.9),
    )
    cap = slide.shapes.add_textbox(Inches(8.1), Inches(6.65), Inches(4.9), Inches(0.4))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Attempt 10 running on every frame (looped, first half of dataset)."
    cp.runs[0].font.size = Pt(12)
    cp.runs[0].font.italic = True
    cp.runs[0].font.color.rgb = TEXT_MUTED

    # Headline numbers strip
    headline = slide.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(7.4), Inches(0.7))
    p = headline.text_frame.paragraphs[0]
    p.text = "Attempt 10: 97.3 % pass   |   30.3 ms mean   |   33 FPS   |   PASS"
    p.runs[0].font.size = Pt(18)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT_GREEN


def slide_hard_dataset(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Hard dataset - FPV / ATV (120 frames)",
              "Ground-level + low altitude FPV. Treelines, moving canopy, occasional sky-only or ground-only frames.")

    add_image(slide, CHART_DIR / "fpv_progression.png",
              Inches(0.4), Inches(1.6), width=Inches(7.4))

    add_video_with_autoplay_loop(
        slide,
        VIDEO_DIR / "fpv_overlay.mp4",
        HERO_DIR / "01_input.png",
        Inches(8.1), Inches(1.7), Inches(4.9), Inches(4.9),
    )
    cap = slide.shapes.add_textbox(Inches(8.1), Inches(6.65), Inches(4.9), Inches(0.4))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Same detector, harder footage."
    cp.runs[0].font.size = Pt(12)
    cp.runs[0].font.italic = True
    cp.runs[0].font.color.rgb = TEXT_MUTED

    headline = slide.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(7.4), Inches(0.7))
    p = headline.text_frame.paragraphs[0]
    p.text = "Attempt 10: 60.0 % pass   |   37.2 ms mean   |   27 FPS   |   speed PASS, accuracy work to do"
    p.runs[0].font.size = Pt(16)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT_GREEN


def slide_algorithm_a(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "How it works  (1/2)  -  see two views of every pixel",
              "Stages 1-3 on one frame from the easy dataset.")

    # 3 columns of stage images, with caption strips beneath each.
    images = [
        (HERO_DIR / "01_input.png", "Input frame", "the raw BGR image."),
        (HERO_DIR / "02b_lab_b.png", "Lab b* channel",
         "Blue sky -> dark, vegetation -> bright. Survives glare and overcast."),
        (HERO_DIR / "03b_otsu_b.png", "Otsu split",
         "Auto threshold on each channel produces a sky/ground mask."),
    ]
    width = Inches(4.0)
    height = Inches(4.4)
    top = Inches(1.7)
    lefts = [Inches(0.4), Inches(4.65), Inches(8.9)]
    for (img, title, caption), left in zip(images, lefts):
        add_image(slide, img, left, top, width=width, height=height)
        cap_box = slide.shapes.add_textbox(left, Inches(6.2), width, Inches(0.9))
        tf = cap_box.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = title
        tp.runs[0].font.size = Pt(15)
        tp.runs[0].font.bold = True
        tp.runs[0].font.color.rgb = ACCENT_BLUE
        sub = tf.add_paragraph()
        sub.text = caption
        sub.runs[0].font.size = Pt(12)
        sub.runs[0].font.color.rgb = TEXT_MUTED

    # Insight strip at bottom
    add_colored_rect(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4), BG_LIGHT)
    note = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
    np_ = note.text_frame.paragraphs[0]
    np_.text = ("Run the same recipe on grayscale Lab L AND Lab b* - two independent hypotheses, "
                "each survives a different failure mode (haze vs glare).")
    np_.runs[0].font.size = Pt(12)
    np_.runs[0].font.italic = True
    np_.runs[0].font.color.rgb = TEXT_DARK


def slide_algorithm_b(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "How it works  (2/2)  -  pool candidates, pick by physics",
              "Stages 4-6: from boundary edges to the winning line.")

    images = [
        (HERO_DIR / "04b_boundary_b.png",
         "Filtered boundary",
         "Sobel orientation filter discards near-vertical edges (tree trunks, frame borders)."),
        (HERO_DIR / "05_ransac_candidates.png",
         "RANSAC pool",
         "Top-K hypotheses per channel. Inlier count alone is unreliable on tricky frames."),
        (HERO_DIR / "06_ettinger_rerank.png",
         "Coherence rerank",
         "Score = colour separation between regions / scatter within. Picks the line that physically separates two coherent regions."),
    ]
    width = Inches(4.0)
    height = Inches(4.4)
    top = Inches(1.7)
    lefts = [Inches(0.4), Inches(4.65), Inches(8.9)]
    for (img, title, caption), left in zip(images, lefts):
        add_image(slide, img, left, top, width=width, height=height)
        cap_box = slide.shapes.add_textbox(left, Inches(6.2), width, Inches(0.9))
        tf = cap_box.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = title
        tp.runs[0].font.size = Pt(15)
        tp.runs[0].font.bold = True
        tp.runs[0].font.color.rgb = ACCENT_BLUE
        sub = tf.add_paragraph()
        sub.text = caption
        sub.runs[0].font.size = Pt(12)
        sub.runs[0].font.color.rgb = TEXT_MUTED

    add_colored_rect(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4), BG_LIGHT)
    note = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
    np_ = note.text_frame.paragraphs[0]
    np_.text = ("Why coherence beat inlier count: attempt 6 (dual-channel + inlier-count selection) regressed. "
                "Attempt 7 fixed the rule -> +13 pp on FPV.")
    np_.runs[0].font.size = Pt(12)
    np_.runs[0].font.italic = True
    np_.runs[0].font.color.rgb = TEXT_DARK


def slide_performance(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "How we tested performance",
              "Pi 5 isn't on the dev machine - so we approximate its CPU budget in Docker.")

    # Left column - the harness
    left_block = [
        ("The harness", True, 22),
        ("Docker container, single CPU core, 3.5 GB RAM, OMP_NUM_THREADS=1.", False, 16),
        ("Same Python wheels and OpenCV the Pi will run.", False, 16),
        ("Conservative single-core model: real Pi 5 will be at least this fast.", False, 16),
        ("", False, 8),
        ("The pass gate", True, 22),
        ("Mean AND p90 latency  <=  67 ms per frame.", False, 16),
        ("That's the >= 15 FPS requirement, applied to the worst 10 % of frames too.", False, 16),
        ("", False, 8),
        ("How we run it", True, 22),
        ("`tools/bench_docker.sh`  ->  full 490-frame eval  ->  per-attempt JSON.", False, 16),
        ("Re-run on every attempt: regressions caught immediately.", False, 16),
    ]
    add_text_block(slide, Inches(0.5), Inches(1.7), Inches(7.0), Inches(5.0), left_block, size=16)

    # Right column - results card for attempt 10
    add_colored_rect(slide, Inches(7.8), Inches(1.7), Inches(5.2), Inches(5.0), BG_LIGHT, ACCENT_GREEN)
    box = slide.shapes.add_textbox(Inches(8.0), Inches(1.85), Inches(4.9), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    def add_para(text, *, size=15, bold=False, color=TEXT_DARK):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text if text else " "
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = bold
        p.runs[0].font.color.rgb = color
        p.space_after = Pt(4)

    add_para("Attempt 10 on the Pi 5 model", size=20, bold=True, color=ACCENT_GREEN)
    add_para("")
    add_para("Easy  (Horizon-UAV, 490 frames)", bold=True, size=16)
    add_para("  mean  30.3 ms     p90  30.7 ms")
    add_para("  ~ 33 FPS    PASS  (gate 67 ms)", color=ACCENT_GREEN, bold=True)
    add_para("")
    add_para("Hard  (FPV / ATV, 120 frames)", bold=True, size=16)
    add_para("  mean  37.2 ms     p90  38.8 ms")
    add_para("  ~ 27 FPS    PASS  (gate 67 ms)", color=ACCENT_GREEN, bold=True)
    add_para("")
    add_para("Headroom on the harder dataset:", size=13, color=TEXT_MUTED)
    add_para("  29 ms unused per frame for the Hailo detector to consume.", size=13, color=TEXT_MUTED)

    add_footer(slide, "Source: full-eval-results-*.json under each attempt dir | Docker spec: docker/Dockerfile + tools/bench_docker.sh")


def slide_final(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Result + budget", "Stage 7 on the same frame, then a Huber refit on the winner's inliers.")

    # Big final hero image
    add_image(slide, HERO_DIR / "07_final.png", Inches(0.4), Inches(1.6), width=Inches(6.8))

    # Right-hand recap card
    add_colored_rect(slide, Inches(7.5), Inches(1.6), Inches(5.5), Inches(5.0), BG_LIGHT, ACCENT_BLUE)
    box = slide.shapes.add_textbox(Inches(7.7), Inches(1.75), Inches(5.2), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    def add_para(text, *, size=16, bold=False, color=TEXT_DARK):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text if text else " "
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = bold
        p.runs[0].font.color.rgb = color
        p.space_after = Pt(4)

    add_para("Attempt 10 - shipped result", size=22, bold=True, color=ACCENT_GREEN)
    add_para("")
    add_para("Easy (Horizon-UAV)", bold=True)
    add_para("  97.3 % pass  |  30.3 ms  |  33 FPS")
    add_para("")
    add_para("Hard (FPV / ATV)", bold=True)
    add_para("  60.0 % pass  |  37.2 ms  |  27 FPS")
    add_para("")
    add_para("Speed budget", bold=True)
    add_para("  67 ms / frame target  ->  comfortable margin on both datasets")
    add_para("")
    add_para("Open problems", size=14, bold=True, color=TEXT_DARK)
    add_para("  Treelines + low ground horizons - worst remaining failure mass.", size=14, color=TEXT_MUTED)
    add_para("  Real Pi 5 benchmark still pending (Docker model is conservative).", size=14, color=TEXT_MUTED)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    set_slide_size_16x9(prs)

    slide_title(prs)
    slide_problem(prs)
    slide_easy_dataset(prs)
    slide_hard_dataset(prs)
    slide_algorithm_a(prs)
    slide_algorithm_b(prs)
    slide_performance(prs)
    slide_final(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
